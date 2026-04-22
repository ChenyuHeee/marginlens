#!/usr/bin/env python3
"""生成开题答辩PPT v5: 精简5分钟版，含2025最新论文数据，改进图表美观度"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, math

IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "img")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── 现代配色方案 ───
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK = RGBColor(0x2D, 0x2D, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
BLUE = RGBColor(0x20, 0x5C, 0xB8)
LIGHT_BLUE = RGBColor(0x42, 0x88, 0xE0)
SKY = RGBColor(0xE3, 0xEF, 0xFD)
RED = RGBColor(0xD6, 0x3B, 0x3B)
LIGHT_RED = RGBColor(0xFD, 0xE8, 0xE8)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
ORANGE = RGBColor(0xE8, 0x8D, 0x13)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)
PURPLE = RGBColor(0x6A, 0x3D, 0xA8)
LIGHT_PURPLE = RGBColor(0xF0, 0xE6, 0xFF)
ACCENT = RGBColor(0x0C, 0x96, 0x9B)  # teal accent


# ═══════════════════════════════════════════════════════
# 工具函数
# ═══════════════════════════════════════════════════════

def add_bg(slide, color=WHITE):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color


def tb(slide, l, t, w, h, text="", sz=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, fn="Microsoft YaHei"):
    """简洁文本框创建"""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.bold = bold; p.font.color.rgb = color; p.font.name = fn; p.alignment = align
    return tf


def ap(tf, text="", sz=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, sp=0, fn="Microsoft YaHei"):
    """添加段落"""
    p = tf.add_paragraph(); p.text = text; p.font.size = Pt(sz)
    p.font.bold = bold; p.font.color.rgb = color; p.font.name = fn; p.alignment = align
    if sp: p.space_before = Pt(sp)
    return p


def rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(1), radius=None):
    """圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill or WHITE
    if border:
        shape.line.color.rgb = border; shape.line.width = bw
    else:
        shape.line.fill.background()
    return shape


def pill(slide, l, t, w, h, fill=BLUE, text="", sz=12, tc=WHITE, bold=True, fn="Microsoft YaHei"):
    """胶囊形标签"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill; shape.line.fill.background()
    # 设置较大圆角
    if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.5  # 最大圆角
    if text:
        tf = shape.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
        p.font.bold = bold; p.font.color.rgb = tc; p.font.name = fn; p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        shape.text_frame.paragraphs[0].space_after = Pt(0)
    return shape


def header(slide, title, sub=""):
    """页面标题栏 - 现代极简风"""
    # 顶部色带
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tb(slide, 0.8, 0.25, 10, 0.6, title, 26, True, BLACK)
    if sub:
        tb(slide, 0.8, 0.75, 10, 0.35, sub, 13, False, GRAY)


def card(slide, l, t, w, h, title="", items=None, bg=SKY, accent=BLUE, tsz=14, isz=12):
    """美观卡片组件"""
    # 阴影效果（浅色底层）
    rect(slide, l + 0.03, t + 0.03, w, h, RGBColor(0xDD, 0xDD, 0xDD))
    # 主卡片
    rect(slide, l, t, w, h, WHITE, RGBColor(0xE0, 0xE0, 0xE0), Pt(1))
    # 顶部色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    if title:
        tb(slide, l + 0.15, t + 0.12, w - 0.3, 0.4, title, tsz, True, accent)
    if items:
        for j, item in enumerate(items):
            tb(slide, l + 0.2, t + 0.55 + j * 0.38, w - 0.35, 0.35, f"• {item}", isz, False, DARK)


def make_table(slide, rows, cols, left, top, width, height):
    return slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table


def set_cell(tbl, r, c, text, sz=12, bold=False, color=BLACK, align=PP_ALIGN.CENTER, fn="Microsoft YaHei"):
    cell = tbl.cell(r, c); cell.text = ""
    p = cell.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.bold = bold; p.font.color.rgb = color; p.font.name = fn; p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def color_row(tbl, row, fill):
    for c in range(len(tbl.columns)):
        cell = tbl.cell(row, c); cell.fill.solid(); cell.fill.fore_color.rgb = fill


def arrow_shape(slide, l, t, w, h, color=BLUE):
    """右箭头形状"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape


def circle_icon(slide, cx, cy, r, fill, text="", sz=14, tc=WHITE, bold=True):
    """圆形图标"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill; shape.line.fill.background()
    if text:
        tf = shape.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.bold = bold
        p.font.color.rgb = tc; p.font.name = "Microsoft YaHei"; p.alignment = PP_ALIGN.CENTER
    return shape


# ═══════════════════════════════════════════════════════
# Slide 0: 封面 (浅色极简)
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# 顶部色带
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

# 左侧装饰色块
deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.08), Inches(0.12), Inches(7.42))
deco.fill.solid(); deco.fill.fore_color.rgb = BLUE; deco.line.fill.background()
deco2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.12), Inches(0.08), Inches(0.05), Inches(7.42))
deco2.fill.solid(); deco2.fill.fore_color.rgb = ACCENT; deco2.line.fill.background()

tb(slide, 1.5, 1.6, 10.5, 1.0, "面向大模型推理加速的", 34, True, DARK, PP_ALIGN.LEFT)
tb(slide, 1.5, 2.5, 10.5, 1.0, "CXL内存扩展与异构计算架构研究", 40, True, BLUE, PP_ALIGN.LEFT)

# 细线分隔
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.7), Inches(3.5), Pt(2))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

tb(slide, 1.5, 4.1, 10, 0.5, "SRTP 开题报告", 20, False, GRAY, PP_ALIGN.LEFT)
tb(slide, 1.5, 5.0, 10, 0.5, "何宸禹    蔡雨禾    黄绎睿", 22, True, DARK, PP_ALIGN.LEFT)
tb(slide, 1.5, 5.6, 10, 0.5, "指导老师： 杨定裕", 16, False, GRAY, PP_ALIGN.LEFT)
tb(slide, 1.5, 6.2, 10, 0.4, "2026年4月", 14, False, LIGHT_GRAY, PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════
# Slide 1: 大模型推理流程 (新增科普页)
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
header(slide, "背景：大模型（LLM）推理是怎么工作的？", "从用户输入到逐字输出——理解推理流程与 KV Cache 的作用")

# ──── 上半: 推理流程图 ────
# Step 1: 用户输入
rect(slide, 0.6, 1.2, 2.2, 1.8, LIGHT_GREEN, GREEN, Pt(2))
circle_icon(slide, 1.0, 1.4, 0.2, GREEN, "1", 12)
tb(slide, 1.3, 1.25, 1.4, 0.3, "用户输入", 14, True, GREEN)
slide.shapes.add_picture(os.path.join(IMG_DIR, "input.png"), Inches(1.1), Inches(1.6), Inches(1.2), Inches(1.2))

# Arrow 1→2
arrow_shape(slide, 2.95, 1.95, 0.5, 0.25, RGBColor(0xCC, 0xCC, 0xCC))

# Step 2: Prefill 阶段
rect(slide, 3.6, 1.2, 3.2, 1.8, SKY, BLUE, Pt(2))
circle_icon(slide, 4.0, 1.4, 0.2, BLUE, "2", 12)
tb(slide, 4.3, 1.25, 2.4, 0.3, "Prefill 阶段", 14, True, BLUE)
slide.shapes.add_picture(os.path.join(IMG_DIR, "prefill.png"), Inches(3.8), Inches(1.6), Inches(1.1), Inches(1.1))
tf = tb(slide, 4.95, 1.65, 1.7, 0.35, "并行处理输入", 12, False, DARK)
ap(tf, "→ 计算密集", 12, False, BLUE, sp=3)
ap(tf, "→ 生成 KV Cache", 12, True, RED, sp=3)

# Arrow 2→3
arrow_shape(slide, 6.95, 1.95, 0.5, 0.25, RGBColor(0xCC, 0xCC, 0xCC))

# Step 3: Decode 阶段
rect(slide, 7.6, 1.2, 3.2, 1.8, LIGHT_ORANGE, ORANGE, Pt(2))
circle_icon(slide, 8.0, 1.4, 0.2, ORANGE, "3", 12)
tb(slide, 8.3, 1.25, 2.4, 0.3, "Decode 阶段", 14, True, ORANGE)
slide.shapes.add_picture(os.path.join(IMG_DIR, "decode.png"), Inches(7.8), Inches(1.6), Inches(1.1), Inches(1.1))
tf = tb(slide, 8.95, 1.65, 1.7, 0.35, "逐token生成", 12, False, DARK)
ap(tf, "→ 内存密集", 12, False, ORANGE, sp=3)
ap(tf, "→ KV Cache 增长", 12, True, RED, sp=3)

# Step 4: 输出
arrow_shape(slide, 10.95, 1.95, 0.5, 0.25, RGBColor(0xCC, 0xCC, 0xCC))
rect(slide, 11.6, 1.2, 1.4, 1.8, LIGHT_PURPLE, PURPLE, Pt(2))
circle_icon(slide, 11.9, 1.4, 0.2, PURPLE, "4", 12)
tb(slide, 11.65, 1.55, 1.3, 0.3, "逐字输出", 13, True, PURPLE, PP_ALIGN.CENTER)
slide.shapes.add_picture(os.path.join(IMG_DIR, "output.png"), Inches(11.7), Inches(1.85), Inches(1.1), Inches(1.1))


# ──── 中: KV Cache 核心解释 ────
rect(slide, 0.6, 3.15, 12.3, 2.2, WHITE, BLUE, Pt(2))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.6), Inches(3.15), Inches(12.3), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

tb(slide, 0.9, 3.25, 5, 0.35, "什么是 KV Cache？为什么它是瓶颈？", 16, True, BLUE)

# KV Cache 插图
slide.shapes.add_picture(os.path.join(IMG_DIR, "kv cache.png"), Inches(0.8), Inches(3.7), Inches(1.3), Inches(1.3))

# 左: 解释
tf = tb(slide, 2.2, 3.7, 4.2, 0.35, "KV Cache = 每层 Transformer 的 Key & Value 缓存", 13, True, DARK)
ap(tf, "", 4)
ap(tf, "• 避免重复计算：每生成 1 个 token 需要回看所有历史 token", 13, False, DARK, sp=5)
ap(tf, "• 没有 KV Cache → 每步都重算 → 极慢（O(n²) → O(n)）", 13, False, DARK, sp=4)
ap(tf, "• 代价：KV Cache 需常驻 GPU 显存，且随序列增长线性膨胀", 13, True, RED, sp=4)

# 右: 可视化 KV Cache 增长
# 画一个简化的 KV Cache 增长条形图
tb(slide, 7.0, 3.3, 5.5, 0.35, "KV Cache 随序列长度线性增长", 14, True, RED, PP_ALIGN.CENTER)

ctx_lens  = ["1K", "4K", "16K", "64K", "128K"]
kv_sizes  = [0.15, 0.6, 2.4, 9.6, 19.2]  # 大致 GB (70B model)
max_kv = 19.2
bar_max_w = 4.8

for i, (cl, kvs) in enumerate(zip(ctx_lens, kv_sizes)):
    by = 3.75 + i * 0.3
    # 标签
    tb(slide, 7.0, by, 0.7, 0.28, cl, 11, True, DARK, PP_ALIGN.RIGHT)
    # 条
    bw = (kvs / max_kv) * bar_max_w
    bar_color = RED if kvs > 5 else (ORANGE if kvs > 1 else GREEN)
    rect(slide, 7.8, by + 0.02, max(bw, 0.15), 0.22, bar_color, bar_color, Pt(0))
    # 数值
    tb(slide, 7.85 + bw, by, 1.2, 0.28, f"{kvs:.1f} GB" if kvs >= 1 else f"{kvs*1024:.0f} MB", 10, True, bar_color)

tb(slide, 7.0, 5.05, 5.5, 0.25, "（以 70B 模型为例，FP16 精度）", 10, False, GRAY, PP_ALIGN.CENTER)

# ──── 下: 核心矛盾总结 ────
rect(slide, 0.6, 5.55, 12.3, 1.3, RGBColor(0xFF, 0xF8, 0xF0), ORANGE, Pt(2))
tb(slide, 0.9, 5.6, 11.7, 0.35, "核心矛盾", 16, True, ORANGE)

# 三列对比
conflict_items = [
    ("GPU HBM 容量", "80 GB（单卡上限）", "有限", RED),
    ("KV Cache 需求", "70B+128K → 160 GB", "远超 HBM", RED),
    ("用户期望", "更长上下文、更多并发", "持续增长", ORANGE),
]
for j, (label, val, trend, clr) in enumerate(conflict_items):
    xp = 0.9 + j * 4.1
    rect(slide, xp, 6.0, 3.8, 0.7, WHITE, clr, Pt(1.5))
    tb(slide, xp + 0.1, 6.0, 3.6, 0.3, label, 13, True, clr, PP_ALIGN.CENTER)
    tb(slide, xp + 0.1, 6.3, 2.5, 0.3, val, 12, False, DARK)
    pill(slide, xp + 2.7, 6.32, 0.9, 0.25, clr, trend, 10)


# ═══════════════════════════════════════════════════════
# Slide 2: 研究背景 – 内存墙 (精简合并)
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
header(slide, "研究背景：大模型推理的\"内存墙\"困境")

# GPU+内存条插图
slide.shapes.add_picture(os.path.join(IMG_DIR, "GPU+内存条.png"), Inches(11.5), Inches(0.15), Inches(1.1), Inches(1.1))

# ──── 左: 剪刀差图 (改进视觉) ────
tb(slide, 0.8, 1.15, 6, 0.35, "GPU 算力 vs 内存增长（相对 V100 = 1×）", 14, True, DARK)

gpus = ["V100\n(2017)", "A100\n(2020)", "H100\n(2022)", "B200\n(2024)"]
x_ps = [1.5, 3.0, 4.5, 6.0]
compute = [1.0, 2.5, 7.9, 14.4]
hbm_cap = [1.0, 2.5, 2.5, 6.0]

max_h = 3.5
base_y = 5.1
# 网格线
for frac in [0.25, 0.5, 0.75, 1.0]:
    gh = max_h * frac
    gl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(1.3), Inches(base_y - gh), Inches(5.8), Pt(0.5))
    gl.fill.solid(); gl.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8); gl.line.fill.background()

for i, (xp, gpu) in enumerate(zip(x_ps, gpus)):
    # 算力柱
    bh = (compute[i] / max(compute)) * max_h
    rect(slide, xp, base_y - bh, 0.5, bh, BLUE, BLUE, Pt(0))
    tb(slide, xp - 0.1, base_y - bh - 0.32, 0.75, 0.3, f"{compute[i]:.0f}×" if compute[i] >= 2 else "1×",
       11, True, BLUE, PP_ALIGN.CENTER)
    # HBM柱
    bh2 = (hbm_cap[i] / max(compute)) * max_h
    rect(slide, xp + 0.6, base_y - bh2, 0.5, bh2, ORANGE, ORANGE, Pt(0))
    tb(slide, xp + 0.5, base_y - bh2 - 0.32, 0.75, 0.3, f"{hbm_cap[i]:.0f}×" if hbm_cap[i] >= 2 else "1×",
       11, True, ORANGE, PP_ALIGN.CENTER)
    # GPU标签
    tb(slide, xp - 0.05, base_y + 0.05, 1.2, 0.5, gpu, 11, True, DARK, PP_ALIGN.CENTER)

# 图例
pill(slide, 1.3, 5.7, 1.2, 0.3, BLUE, "FP16 算力", 11)
pill(slide, 2.7, 5.7, 1.2, 0.3, ORANGE, "HBM 容量", 11)

# 剪刀差标注
rect(slide, 0.8, 6.15, 7.0, 0.55, LIGHT_RED, RED, Pt(2))
tb(slide, 1.0, 6.2, 6.6, 0.4, "⚠ 算力增长 14× vs HBM容量仅 6× → Gap持续扩大", 13, True, RED, PP_ALIGN.CENTER)

# 右上: 大字强调 gap 数据
rect(slide, 8.0, 6.15, 4.8, 0.55, RGBColor(0xFF, 0xF0, 0xF0), RED, Pt(2))
tf = tb(slide, 8.15, 6.15, 4.5, 0.5, "", 13)
p = ap(tf, "→ GPU HBM (80GB) 远远不够！", 14, True, RED, PP_ALIGN.CENTER, sp=0)

# ──── 右: KV Cache 问题 ────
tb(slide, 8.0, 1.15, 5, 0.35, "KV Cache：内存瓶颈核心", 14, True, DARK)

# 模型表
tbl = make_table(slide, 5, 3, 8.0, 1.6, 4.8, 2.0)
for j, h in enumerate(["模型", "权重(FP16)", "KV / token"]):
    set_cell(tbl, 0, j, h, 12, True, WHITE)
color_row(tbl, 0, BLUE)
for i, (m, w, kv) in enumerate([
    ("LLaMA-3-8B", "16 GB", "~0.5 MB"),
    ("LLaMA-3-70B", "140 GB", "~2.5 MB"),
    ("Qwen-32B", "64 GB", "~1.3 MB"),
    ("DeepSeek-V3", "~340 GB(MoE)", "-"),
]):
    set_cell(tbl, i+1, 0, m, 12, True); set_cell(tbl, i+1, 1, w, 12)
    bold_kv = "MB" in kv
    set_cell(tbl, i+1, 2, kv, 12, bold_kv, RED if bold_kv else BLACK)
    if i % 2 == 0: color_row(tbl, i+1, SKY)

# KV cache 数据点 - 改为卡片式布局更醒目
rect(slide, 8.0, 3.85, 4.8, 2.85, LIGHT_RED, RED, Pt(1.5))
tf = tb(slide, 8.15, 3.9, 4.5, 0.35, "KV Cache 压力数据", 14, True, RED)

# 核心数据用大字高亮
rect(slide, 8.2, 4.35, 2.15, 0.85, WHITE, RED, Pt(1.5))
tb(slide, 8.25, 4.35, 2.05, 0.3, "70B + 128K ctx", 11, False, GRAY, PP_ALIGN.CENTER)
tb(slide, 8.25, 4.65, 2.05, 0.45, "~160 GB", 22, True, RED, PP_ALIGN.CENTER)

rect(slide, 10.55, 4.35, 2.15, 0.85, WHITE, ORANGE, Pt(1.5))
tb(slide, 10.6, 4.35, 2.05, 0.3, "Kimi 50M token", 11, False, GRAY, PP_ALIGN.CENTER)
tb(slide, 10.6, 4.65, 2.05, 0.45, "~20 TB", 22, True, ORANGE, PP_ALIGN.CENTER)

# 补充数据
tf = tb(slide, 8.15, 5.3, 4.5, 0.35, "", 13)
ap(tf, "• KV 占 GPU 显存 > 30% [vLLM]", 13, False, DARK, sp=4)
ap(tf, "• HBM 14秒内被 KV 占满 [CachedAttention]", 13, False, DARK, sp=4)


# ═══════════════════════════════════════════════════════
# Slide 2: P/D解耦与CXL机遇 (精简合并)
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
header(slide, "Prefill-Decode 解耦趋势与 CXL 机遇")

# CXL芯片插图
slide.shapes.add_picture(os.path.join(IMG_DIR, "CXL.png"), Inches(11.5), Inches(0.15), Inches(1.1), Inches(1.1))

# ──── 上: P/D 解耦架构图 ────
# Prefill GPU
rect(slide, 1.0, 1.3, 2.2, 1.5, LIGHT_GREEN, GREEN, Pt(2))
tb(slide, 1.2, 1.35, 1.8, 0.35, "Prefill GPU", 14, True, GREEN, PP_ALIGN.CENTER)
tf = tb(slide, 1.15, 1.7, 1.9, 0.3, "计算密集型", 12, False, DARK, PP_ALIGN.CENTER)
ap(tf, "处理输入 prompt", 11, False, GRAY, PP_ALIGN.CENTER, sp=2)
ap(tf, "生成 KV Cache", 11, True, GREEN, PP_ALIGN.CENTER, sp=2)

# 箭头1: Prefill → KV Transfer
arrow_shape(slide, 3.4, 1.85, 0.8, 0.3, RGBColor(0xCC, 0xCC, 0xCC))
tb(slide, 3.35, 1.5, 0.9, 0.3, "KV传输", 10, True, RED, PP_ALIGN.CENTER)
tb(slide, 3.35, 2.2, 0.9, 0.25, "瓶颈!", 11, True, RED, PP_ALIGN.CENTER)

# Decode GPU
rect(slide, 4.4, 1.3, 2.2, 1.5, LIGHT_ORANGE, ORANGE, Pt(2))
tb(slide, 4.6, 1.35, 1.8, 0.35, "Decode GPU", 14, True, ORANGE, PP_ALIGN.CENTER)
tf = tb(slide, 4.55, 1.7, 1.9, 0.3, "内存密集型", 12, False, DARK, PP_ALIGN.CENTER)
ap(tf, "逐 token 生成", 11, False, GRAY, PP_ALIGN.CENTER, sp=2)
ap(tf, "读取 KV Cache", 11, True, ORANGE, PP_ALIGN.CENTER, sp=2)

# 数据佐证
rect(slide, 7.0, 1.2, 5.7, 1.7, SKY, BLUE, Pt(1))
tf = tb(slide, 7.15, 1.25, 5.4, 0.35, "P/D 解耦已成主流", 14, True, BLUE)
ap(tf, "• DistServe: goodput 提升 7.4× [OSDI'24]", 13, False, DARK, sp=4)
ap(tf, "• Splitwise: 同成本吞吐 1.4× [ISCA'24]", 13, False, DARK, sp=3)
ap(tf, "• NVIDIA Dynamo: 生产级 P/D 框架 [2025]", 13, False, DARK, sp=3)
ap(tf, "⚡ 但 KV 传输依赖 RDMA，延迟高、编程复杂", 13, True, RED, sp=4)

# ──── 下: CXL 内存层次 (横向流程) ────
tb(slide, 0.8, 3.2, 12, 0.4, "CXL：填补内存层次鸿沟", 16, True, DARK)

tiers_data = [
    ("GPU HBM", "80GB\n3.35 TB/s\n~10 ns", RGBColor(0xE8, 0x56, 0x56), 0.8),
    ("CPU DRAM", "2 TB\n200 GB/s\n~90 ns", RGBColor(0xF0, 0xA0, 0x30), 3.9),
    ("CXL Memory", "TB级\n38 GB/s\n~200 ns", BLUE, 7.0),
    ("NVMe SSD", "数TB\n7 GB/s\n~10 μs", RGBColor(0x88, 0x88, 0x88), 10.1),
]
for i, (name, spec, clr, xp) in enumerate(tiers_data):
    is_cxl = (i == 2)
    bw_pt = Pt(3) if is_cxl else Pt(1)
    rect(slide, xp, 3.7, 2.8, 1.5, WHITE, clr, bw_pt)
    # 顶部色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(xp), Inches(3.7), Inches(2.8), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = clr; bar.line.fill.background()
    tb(slide, xp + 0.1, 3.82, 2.6, 0.35, name, 14, True, clr, PP_ALIGN.CENTER)
    tb(slide, xp + 0.1, 4.2, 2.6, 0.8, spec, 12, False, GRAY, PP_ALIGN.CENTER)
    if is_cxl:
        pill(slide, xp + 0.5, 5.0, 1.8, 0.25, BLUE, "← 新增层", 11, WHITE)
    # 箭头
    if i < 3:
        next_xp = tiers_data[i+1][3]
        ax = xp + 2.8 + 0.1
        aw = next_xp - ax - 0.1
        if aw > 0:
            arrow_shape(slide, ax, 4.25, aw, 0.25, RGBColor(0xCC, 0xCC, 0xCC))

# 底部核心优势 - 三张卡片
advantages = [
    ("容量扩展", "单节点 → TB级\n满足长上下文 KV Cache", BLUE, "TB"),
    ("低延迟访问", "load/store语义, ~200ns\n比RDMA快 6.2×", ACCENT, "6×"),
    ("内存池化", "CXL 2.0 多主机共享\n可回收25%搁浅内存", GREEN, "25%"),
]
for j, (title, desc, clr, big_num) in enumerate(advantages):
    xp = 0.8 + j * 4.0
    rect(slide, xp, 5.5, 3.7, 1.5, WHITE, clr, Pt(2))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(xp), Inches(5.5), Inches(3.7), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = clr; bar.line.fill.background()
    circle_icon(slide, xp + 0.45, 6.0, 0.35, clr, big_num, 14)
    tb(slide, xp + 0.95, 5.6, 2.6, 0.35, title, 14, True, clr)
    tb(slide, xp + 0.95, 5.95, 2.6, 0.6, desc, 11, False, GRAY)


# ═══════════════════════════════════════════════════════
# Slide 3: 最新研究验证 – TraCT & Beluga (★ NEW)
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
header(slide, "最新研究验证：CXL 显著提升 LLM 推理性能", "2025年最新论文实证——CXL + LLM 是高度活跃的前沿方向")

# ──── 左: TraCT ────
rect(slide, 0.6, 1.25, 6.0, 5.5, WHITE, BLUE, Pt(2))
# 顶部色带
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.6), Inches(1.25), Inches(6.0), Inches(0.55))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
tb(slide, 0.8, 1.28, 5.6, 0.45, "TraCT [arXiv'25/12, Virginia Tech + SK Hynix]", 15, True, WHITE)

tf = tb(slide, 0.8, 1.95, 5.6, 0.4, "CXL共享内存 替代 RDMA 进行 P/D 间 KV 传输", 13, True, BLUE)
ap(tf, "", 4)
ap(tf, "核心思想", 14, True, DARK, sp=6)
ap(tf, "• 用 CXL 共享内存同时做 KV 传输 + prefix cache", 13, False, DARK, sp=4)
ap(tf, "• GPU 通过 DMA 直接读写 CXL 内存, 消除 NIC", 13, False, DARK, sp=3)
ap(tf, "• 基于 NVIDIA Dynamo 框架实现", 13, False, DARK, sp=3)

# 结果高亮
rect(slide, 0.9, 4.0, 5.4, 2.5, RGBColor(0xF0, 0xF8, 0xFF), BLUE, Pt(1))
tb(slide, 1.1, 4.05, 5.0, 0.35, "关键结果 (vs RDMA-based)", 14, True, BLUE)

metrics_t = [("TTFT 降低", "9.8×", BLUE), ("P99 延迟降低", "6.2×", ACCENT), ("吞吐提升", "1.6×", GREEN)]
for j, (label, val, clr) in enumerate(metrics_t):
    bx = 1.1 + j * 1.8
    rect(slide, bx, 4.45, 1.6, 1.0, WHITE, clr, Pt(1.5))
    tb(slide, bx + 0.05, 4.5, 1.5, 0.3, label, 12, False, GRAY, PP_ALIGN.CENTER)
    tb(slide, bx + 0.05, 4.8, 1.5, 0.5, val, 24, True, clr, PP_ALIGN.CENTER)

tb(slide, 1.1, 5.55, 5.0, 0.35, "→ 对应本课题 方向一", 14, True, BLUE)

# ──── 右: Beluga ────
rect(slide, 6.9, 1.25, 6.0, 5.5, WHITE, GREEN, Pt(2))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(6.9), Inches(1.25), Inches(6.0), Inches(0.55))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()
tb(slide, 7.1, 1.28, 5.6, 0.45, "Beluga [SIGMOD'26, 阿里巴巴]", 15, True, WHITE)

tf = tb(slide, 7.1, 1.95, 5.6, 0.4, "CXL Switch 内存池 管理 LLM KVCache", 13, True, GREEN)
ap(tf, "", 4)
ap(tf, "核心思想", 14, True, DARK, sp=6)
ap(tf, "• CXL 2.0 Switch 构建 8TB 共享内存池", 13, False, DARK, sp=4)
ap(tf, "• GPU/CPU 均可直接访问 CXL 内存池", 13, False, DARK, sp=3)
ap(tf, "• 集成到 vLLM, 简化调度 (无需考虑KV局部性)", 13, False, DARK, sp=3)

# 结果高亮
rect(slide, 7.2, 4.0, 5.4, 2.5, RGBColor(0xF0, 0xFF, 0xF0), GREEN, Pt(1))
tb(slide, 7.4, 4.05, 5.0, 0.35, "关键结果 (vs RDMA/MoonCake)", 14, True, GREEN)

metrics_b = [("TTFT 降低", "89.6%", GREEN), ("吞吐提升", "7.35×", ACCENT), ("写延迟降低", "7.0×", BLUE)]
for j, (label, val, clr) in enumerate(metrics_b):
    bx = 7.4 + j * 1.8
    rect(slide, bx, 4.45, 1.6, 1.0, WHITE, clr, Pt(1.5))
    tb(slide, bx + 0.05, 4.5, 1.5, 0.3, label, 12, False, GRAY, PP_ALIGN.CENTER)
    tb(slide, bx + 0.05, 4.8, 1.5, 0.5, val, 24, True, clr, PP_ALIGN.CENTER)

tb(slide, 7.4, 5.55, 5.0, 0.35, "→ 对应本课题 方向一 + 方向二", 14, True, GREEN)

# 底部总结
rect(slide, 0.6, 6.85, 12.3, 0.45, BLACK)
tb(slide, 0.8, 6.87, 11.9, 0.4,
   "✦ CXL + LLM 推理 已被顶会验证可行——本课题恰好处于这一前沿交叉点", 14, True, LIGHT_BLUE, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# Slide 4: 两大研究方向 (精华总览)
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
header(slide, "两大研究方向")

# ──── 方向一 (左) ────
rect(slide, 0.6, 1.1, 6.0, 5.7, WHITE, BLUE, Pt(2))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.6), Inches(1.1), Inches(6.0), Inches(0.6))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

circle_icon(slide, 1.15, 1.4, 0.3, WHITE, "1", 18, BLUE, True)
tb(slide, 1.6, 1.18, 4.8, 0.5, "Serverless 场景下 P/D 间\nKV Cache 通讯", 16, True, WHITE)

# 方向一插图: 服务器+共享内存池
slide.shapes.add_picture(os.path.join(IMG_DIR, "memory pool.png"), Inches(4.8), Inches(1.15), Inches(0.55), Inches(0.55))

# 问题
tf = tb(slide, 0.8, 1.9, 5.6, 0.35, "问题", 14, True, RED)
ap(tf, "P/D 解耦后 KV 传输依赖 RDMA, 延迟高; Serverless", 13, False, DARK, sp=4)
ap(tf, "实例弹性伸缩, KV Cache 无处存放", 13, False, DARK, sp=2)

# 方案
tf = tb(slide, 0.8, 3.0, 5.6, 0.35, "CXL 方案", 14, True, BLUE)
ap(tf, "• CXL 共享内存池实现 P/D 间零拷贝 KV 传递", 13, False, DARK, sp=4)
ap(tf, "• load/store 语义, 比 RDMA 快 6.2× [DirectCXL]", 13, True, BLUE, sp=3)
ap(tf, "• KV Cache 独立于实例生命周期 (存活在 CXL 池)", 13, False, DARK, sp=3)

# 研究内容 - 加圆形编号
tb(slide, 0.8, 4.5, 5.6, 0.35, "研究内容", 14, True, DARK)
research_1 = [
    ("①", "CXL KV Cache 共享协议设计", BLUE),
    ("②", "Serverless 弹性 KV 生命周期管理", ACCENT),
    ("③", "与 RDMA/TCP 方案的对比评测", GREEN),
]
for k, (num, text, clr) in enumerate(research_1):
    ry = 4.95 + k * 0.38
    circle_icon(slide, 1.1, ry + 0.12, 0.15, clr, num, 11)
    tb(slide, 1.4, ry - 0.02, 4.8, 0.35, text, 13, False, DARK)

# 参考
pill(slide, 0.8, 6.15, 5.6, 0.35, SKY, "相关: TraCT · DistServe · Mooncake · Splitwise · Pond", 11, BLUE, False)

# ──── 方向二 (右) ────
rect(slide, 6.9, 1.1, 6.0, 5.7, WHITE, GREEN, Pt(2))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(6.9), Inches(1.1), Inches(6.0), Inches(0.6))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()

circle_icon(slide, 7.45, 1.4, 0.3, WHITE, "2", 18, GREEN, True)
tb(slide, 7.9, 1.18, 4.8, 0.5, "KV Cache Offloading\n到 CXL 内存", 16, True, WHITE)

# 方向二插图: 多层存储金字塔
slide.shapes.add_picture(os.path.join(IMG_DIR, "pyramid.png"), Inches(11.1), Inches(1.15), Inches(0.55), Inches(0.55))

# 问题
tf = tb(slide, 7.1, 1.9, 5.6, 0.35, "问题", 14, True, RED)
ap(tf, "GPU HBM (80GB) 远不够; 现有 offloading 带宽低:", 13, False, DARK, sp=4)
ap(tf, "GPU↔CPU 仅 12 GB/s, CPU↔SSD 仅 2 GB/s", 13, True, RED, sp=2)

# 方案
tf = tb(slide, 7.1, 3.0, 5.6, 0.35, "CXL 方案", 14, True, GREEN)
ap(tf, "• 4 层架构: HBM → DRAM → CXL → SSD", 13, True, GREEN, sp=4)
ap(tf, "• CXL: TB级容量, ~200ns 延迟, 38 GB/s 带宽", 13, False, DARK, sp=3)
ap(tf, "• TPP 透明分层: 仅 20% 本地 DRAM 性能差 <1%", 13, False, DARK, sp=3)

# 研究内容 - 加圆形编号
tb(slide, 7.1, 4.5, 5.6, 0.35, "研究内容", 14, True, DARK)
research_2 = [
    ("①", "CXL-aware KV 热度感知分层策略", GREEN),
    ("②", "vLLM/SGLang CXL 后端集成", BLUE),
    ("③", "多模型/多上下文长度 benchmark", ACCENT),
]
for k, (num, text, clr) in enumerate(research_2):
    ry = 4.95 + k * 0.38
    circle_icon(slide, 7.4, ry + 0.12, 0.15, clr, num, 11)
    tb(slide, 7.7, ry - 0.02, 4.8, 0.35, text, 13, False, DARK)

# 参考
pill(slide, 7.1, 6.15, 5.6, 0.35, LIGHT_GREEN, "相关: Beluga · CachedAttention · FlexGen · HeteGen · TPP", 11, GREEN, False)

# 底部连接
rect(slide, 0.6, 6.9, 12.3, 0.4, BLACK)
tb(slide, 0.8, 6.92, 11.9, 0.35,
   "两个方向互补：方向一解决 P/D 间 KV 传输效率 → 方向二解决单实例 KV 容量瓶颈", 13, True, LIGHT_BLUE, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# Slide 5: 方向一详细架构
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
header(slide, "方向一：CXL 共享内存替代 RDMA 进行 KV 传输")

# ──── 上半部分: VS 对比布局 ────
# 中央 VS 分割标志
pill(slide, 6.15, 2.3, 1.0, 0.55, RED, "VS", 18, WHITE)

# ── 左: RDMA 当前方案 ──
rect(slide, 0.5, 1.15, 5.5, 3.6, RGBColor(0xFB, 0xFB, 0xFB), RGBColor(0xDD, 0xDD, 0xDD), Pt(1))
# 标题区
rect(slide, 0.5, 1.15, 5.5, 0.55, RGBColor(0xF0, 0xF0, 0xF0), RGBColor(0xDD, 0xDD, 0xDD), Pt(1))
tb(slide, 0.7, 1.2, 5.0, 0.45, "❌  RDMA KV 传输 (当前)", 16, True, GRAY, PP_ALIGN.CENTER)

# 简化流程: 3步传输路径
flow_rdma = [
    (0.8, 1.95, 1.5, 0.7, "Prefill GPU", LIGHT_GREEN, GREEN),
    (2.8, 1.95, 1.5, 0.7, "RDMA\nNIC × 2", RGBColor(0xF5, 0xF0, 0xF0), RGBColor(0xCC, 0x88, 0x88)),
    (4.8, 1.95, 1.5, 0.7, "Decode GPU", LIGHT_ORANGE, ORANGE),
]
for (bx, by, bw, bh, bt, bg_c, bd_c) in flow_rdma:
    rect(slide, bx, by, bw, bh, bg_c, bd_c, Pt(2))
    tb(slide, bx, by + 0.1, bw, bh - 0.2, bt, 12, True, bd_c, PP_ALIGN.CENTER)
# 流程箭头
arrow_shape(slide, 2.38, 2.2, 0.35, 0.2, RGBColor(0xCC, 0xCC, 0xCC))
arrow_shape(slide, 4.38, 2.2, 0.35, 0.2, RGBColor(0xCC, 0xCC, 0xCC))

# 路径说明
tb(slide, 0.8, 2.75, 5.0, 0.3, "GPU → DRAM → NIC → NIC → DRAM → GPU", 11, False, GRAY, PP_ALIGN.CENTER)

# 问题列表
problems = [
    ("6 次内存拷贝", "数据搬运开销巨大"),
    ("75% 同步延迟", "RDMA 同步协议瓶颈 [Beluga]"),
    ("编程复杂", "QP 管理 + sglists 限制"),
]
for i, (title, desc) in enumerate(problems):
    y = 3.15 + i * 0.55
    # 红色圆点
    circle_icon(slide, 0.85, y + 0.05, 0.28, RED, str(i + 1), 10)
    tb(slide, 1.2, y, 2.0, 0.3, title, 12, True, RED)
    tb(slide, 3.1, y, 2.8, 0.3, desc, 11, False, GRAY)

# ── 右: CXL 方案 ──
rect(slide, 7.3, 1.15, 5.5, 3.6, RGBColor(0xF5, 0xF9, 0xFF), BLUE, Pt(2))
# 标题区
rect(slide, 7.3, 1.15, 5.5, 0.55, SKY, BLUE, Pt(2))
tb(slide, 7.5, 1.2, 5.1, 0.45, "✓  CXL 共享内存 (本方案)", 16, True, BLUE, PP_ALIGN.CENTER)

# 简化流程: 3节点
flow_cxl = [
    (7.6, 1.95, 1.4, 0.7, "Prefill GPU", LIGHT_GREEN, GREEN),
    (9.5, 1.95, 1.8, 0.7, "CXL\n共享内存池", SKY, BLUE),
    (11.8, 1.95, 1.4, 0.7, "Decode GPU", LIGHT_ORANGE, ORANGE),
]
for (bx, by, bw, bh, bt, bg_c, bd_c) in flow_cxl:
    rect(slide, bx, by, bw, bh, bg_c, bd_c, Pt(2))
    tb(slide, bx, by + 0.1, bw, bh - 0.2, bt, 12, True, bd_c, PP_ALIGN.CENTER)
# 双向箭头
arrow_shape(slide, 9.08, 2.2, 0.35, 0.2, BLUE)
arrow_shape(slide, 11.38, 2.2, 0.35, 0.2, BLUE)

# 路径说明
tb(slide, 7.5, 2.75, 5.1, 0.3, "GPU ↔ CXL Pool (load/store 直接访问)", 11, False, BLUE, PP_ALIGN.CENTER)

# 优势列表
advantages = [
    ("零拷贝传输", "GPU DMA 直达 CXL 内存"),
    ("简单编程", "load/store 语义, 无需 RDMA"),
    ("低延迟", "TTFT ↓9.8× [TraCT 2025]"),
]
for i, (title, desc) in enumerate(advantages):
    y = 3.15 + i * 0.55
    circle_icon(slide, 7.65, y + 0.05, 0.28, GREEN, "✓", 10)
    tb(slide, 8.0, y, 2.0, 0.3, title, 12, True, GREEN)
    tb(slide, 9.9, y, 2.8, 0.3, desc, 11, False, DARK)

# ──── 中间: 延迟对比亮点条 ────
rect(slide, 0.5, 4.9, 12.3, 0.55, SKY, BLUE, Pt(1))
tb(slide, 0.8, 4.95, 2.0, 0.4, "延迟对比 (TTFT)", 14, True, BLUE)
# RDMA
tb(slide, 3.0, 4.97, 1.0, 0.35, "RDMA", 12, True, GRAY, PP_ALIGN.RIGHT)
rect(slide, 4.1, 5.02, 3.3, 0.25, RGBColor(0xDD, 0xDD, 0xDD), RGBColor(0xBB, 0xBB, 0xBB), Pt(0))
tb(slide, 7.45, 4.97, 0.6, 0.35, "9.8×", 12, True, RED)
# CXL
tb(slide, 8.2, 4.97, 1.0, 0.35, "CXL", 12, True, BLUE, PP_ALIGN.RIGHT)
rect(slide, 9.3, 5.02, 0.34, 0.25, BLUE, BLUE, Pt(0))
tb(slide, 9.7, 4.97, 0.5, 0.35, "1×", 12, True, BLUE)
pill(slide, 10.4, 4.98, 2.2, 0.32, GREEN, "↓ 9.8× 更低延迟", 11)

# ──── 下: 研究计划 ────
tb(slide, 0.6, 5.6, 12, 0.35, "研究计划", 15, True, DARK)

tasks_1 = [
    ("1.1 CXL KV共享协议", "P/D 间 KV 寻址、\n一致性保障、\n并发访问控制", BLUE),
    ("1.2 Serverless弹性调度", "KV Cache 独立生命周期\n热度感知缓存淘汰\n实例启停时 KV 迁移", ACCENT),
    ("1.3 性能优化与评测", "预取策略设计\n批量 vs 流式传输\nvs RDMA/TCP 对比", GREEN),
]
for j, (title, desc, clr) in enumerate(tasks_1):
    xp = 0.6 + j * 4.2
    card(slide, xp, 5.95, 3.9, 1.4, title, None, bg=WHITE, accent=clr, tsz=13)
    tb(slide, xp + 0.15, 6.4, 3.6, 0.85, desc, 11, False, DARK)


# ═══════════════════════════════════════════════════════
# Slide 6: 方向二详细架构
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
header(slide, "方向二：KV Cache Offloading 到 CXL 内存")

# ──── 左侧: 分层存储可视化 (水平条形，宽度递增表示容量) ────
tb(slide, 0.6, 1.1, 5.0, 0.4, "四级存储层次", 16, True, DARK)

tiers = [
    ("GPU HBM",    "80 GB | 3.35 TB/s | ~10 ns",   RGBColor(0xE8, 0x56, 0x56), 3.2, "热 KV + 权重"),
    ("CPU DRAM",   "2 TB  | 200 GB/s  | ~90 ns",    ORANGE, 4.2, "温 KV + Buffer"),
    ("CXL Memory", "TB级  | 38 GB/s   | ~200 ns",   BLUE, 5.2, "冷 KV + 复用 ★"),
    ("NVMe SSD",   "数TB  | 7 GB/s    | ~10 μs",    RGBColor(0x88, 0x88, 0x88), 6.2, "历史 KV"),
]
for i, (name, spec, clr, w, usage) in enumerate(tiers):
    y = 1.65 + i * 0.85
    is_cxl = (i == 2)
    bw_pt = Pt(3) if is_cxl else Pt(1)
    bc = clr if is_cxl else RGBColor(0xDD, 0xDD, 0xDD)
    bg_fill = SKY if is_cxl else WHITE

    # 水平条 - 宽度递增表示容量增大
    rect(slide, 0.6, y, w, 0.65, bg_fill, bc, bw_pt)
    # 左侧色块标识
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(y), Inches(0.1), Inches(0.65))
    bar.fill.solid(); bar.fill.fore_color.rgb = clr; bar.line.fill.background()

    # Tier名称
    tb(slide, 0.8, y + 0.02, 1.8, 0.3, name, 13, True, clr)
    # 规格
    tb(slide, 0.8, y + 0.32, w - 0.3, 0.3, spec, 10, False, GRAY)
    # 用途标注（条形右侧）
    tb(slide, w + 0.75, y + 0.12, 2.5, 0.35, usage, 12, is_cxl, BLUE if is_cxl else DARK)

    # 层间热冷箭头
    if i < 3:
        arrow_shape(slide, 1.6, y + 0.65, 0.2, 0.18, RGBColor(0xCC, 0xCC, 0xCC))

# 热→冷标注
tb(slide, 0.1, 1.65, 0.5, 0.3, "热", 12, True, RED, PP_ALIGN.CENTER)
tb(slide, 0.1, 4.05, 0.5, 0.3, "冷", 12, True, GRAY, PP_ALIGN.CENTER)
# 竖线
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.3), Inches(1.95), Inches(0.03), Inches(2.15))
bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD); bar.line.fill.background()

# ──── 右侧上: 问题卡片 ────
rect(slide, 7.8, 1.1, 5.0, 1.8, LIGHT_RED, RED, Pt(2))
# 标题
rect(slide, 7.8, 1.1, 5.0, 0.45, RED, RED, Pt(0))
tb(slide, 8.0, 1.13, 4.6, 0.4, "🚫  现有 Offloading 瓶颈", 15, True, WHITE, PP_ALIGN.CENTER)

items_problem = [
    "FlexGen: GPU↔CPU 仅 12 GB/s",
    "HeteGen: PCIe 仅 25-30 GB/s",
    "SSD 延迟 ~10 μs (100× DRAM)",
    "缺少 DRAM 与 SSD 间的中间层!",
]
for i, itm in enumerate(items_problem):
    y = 1.65 + i * 0.29
    is_last = (i == len(items_problem) - 1)
    tb(slide, 8.1, y, 4.5, 0.28, ("⚠ " if is_last else "• ") + itm, 12, is_last, RED if is_last else DARK)

# ──── 右侧下: 解决方案 + 带宽对比 ────
rect(slide, 7.8, 3.1, 5.0, 1.7, LIGHT_GREEN, GREEN, Pt(2))
rect(slide, 7.8, 3.1, 5.0, 0.45, GREEN, GREEN, Pt(0))
tb(slide, 8.0, 3.13, 4.6, 0.4, "✅  CXL 补齐鸿沟", 15, True, WHITE, PP_ALIGN.CENTER)

# 带宽对比条形图 - 更紧凑
bw_data = [
    ("SSD",  0.5,  "7 GB/s",  RGBColor(0xBB, 0xBB, 0xBB), GRAY),
    ("PCIe", 0.86, "12 GB/s", ORANGE, ORANGE),
    ("CXL",  2.7,  "38 GB/s", GREEN, GREEN),
]
for i, (label, bar_w, val, bar_clr, txt_clr) in enumerate(bw_data):
    y = 3.7 + i * 0.32
    tb(slide, 8.0, y, 1.2, 0.28, label, 11, True, txt_clr, PP_ALIGN.RIGHT)
    rect(slide, 9.3, y + 0.04, bar_w, 0.2, bar_clr, bar_clr, Pt(0))
    tb(slide, 9.3 + bar_w + 0.1, y, 1.0, 0.28, val, 11, True, txt_clr)

pill(slide, 9.0, 4.5, 3.5, 0.28, GREEN, "CXL 带宽 3× PCIe, 5× SSD", 11)

# ──── 下: 研究计划 ────
tb(slide, 0.6, 5.0, 12, 0.35, "研究计划", 15, True, DARK)

tasks_2 = [
    ("2.1 CXL-aware 分层策略", "热度感知: 热KV→HBM\n冷KV→CXL, 预取机制\n扩展 CachedAttention", GREEN),
    ("2.2 推理框架集成", "vLLM/SGLang CXL后端\nPagedAttention CXL块管理\nMMAP/DAX 接口适配", BLUE),
    ("2.3 性能评测优化", "真实CXL + NUMA模拟\n多模型/多ctx benchmark\n目标: 吞吐 ≥ 30%↑", ACCENT),
]
for j, (title, desc, clr) in enumerate(tasks_2):
    xp = 0.6 + j * 4.2
    card(slide, xp, 5.35, 3.9, 2.0, title, None, bg=WHITE, accent=clr, tsz=13)
    tb(slide, xp + 0.15, 5.85, 3.6, 1.3, desc, 11, False, DARK)


# ═══════════════════════════════════════════════════════
# Slide 7: 时间规划 + 预期贡献 (合并)
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
header(slide, "时间规划与预期贡献")

# ──── 左: 时间规划甘特图风格 ────
tb(slide, 0.6, 1.05, 6, 0.4, "时间规划", 16, True, DARK)

phases = [
    ("阶段一\n(1-3月)", "文献调研 + 环境搭建", "CXL模拟环境 · LLM访存profiling · KV热度分析", RGBColor(0x66, 0xBB, 0x6A)),
    ("阶段二\n(3-6月)", "方向一: P/D间CXL KV通讯", "共享协议设计 · Serverless弹性KV · vs RDMA评测", BLUE),
    ("阶段三\n(3-6月)", "方向二: KV Offloading到CXL", "4层分层管理 · 集成vLLM/SGLang · 多模型评测", ACCENT),
    ("阶段四\n(2-3月)", "整合与结题", "联合优化 · 论文撰写 · 开源发布", ORANGE),
]
for i, (phase, title, detail, clr) in enumerate(phases):
    y = 1.55 + i * 1.15
    # 阶段标签
    rect(slide, 0.6, y, 1.4, 0.95, clr)
    tb(slide, 0.65, y + 0.15, 1.3, 0.65, phase, 12, True, WHITE, PP_ALIGN.CENTER)
    # 内容条
    rect(slide, 2.1, y, 5.2, 0.95, WHITE, clr, Pt(1.5))
    tb(slide, 2.25, y + 0.08, 4.9, 0.35, title, 14, True, DARK)
    tb(slide, 2.25, y + 0.48, 4.9, 0.4, detail, 12, False, GRAY)
    # 阶段间连接箭头
    if i < len(phases) - 1:
        arrow_shape(slide, 1.15, y + 0.95 + 0.02, 0.25, 0.14, clr)

# ──── 右: 预期贡献 ────
tb(slide, 7.8, 1.05, 5, 0.4, "预期贡献", 16, True, DARK)

contribs = [
    ("01", "CXL P/D KV 通讯方案", "基于 CXL 共享内存的\nP/D 间零拷贝 KV 传递", BLUE),
    ("02", "CXL KV Offloading 系统", "4 层 KV Cache 分层管理\n集成到 vLLM/SGLang", GREEN),
    ("03", "实际系统验证与落地", "对接蚂蚁智算平台需求\n(CCF-蚂蚁科研基金项目)", ACCENT),
    ("04", "开源系统原型", "可在真实 CXL / 模拟器\n上运行的推理系统原型", ORANGE),
]
for i, (num, title, desc, clr) in enumerate(contribs):
    y = 1.55 + i * 1.15
    # 卡片背景
    rect(slide, 7.8, y, 5.0, 0.95, WHITE, clr, Pt(1.5))
    # 左色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(7.8), Inches(y), Inches(0.06), Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = clr; bar.line.fill.background()
    # 编号圆
    circle_icon(slide, 8.3, y + 0.47, 0.3, clr, num, 16)
    # 内容
    tb(slide, 8.7, y + 0.05, 4.0, 0.35, title, 15, True, DARK)
    tb(slide, 8.7, y + 0.4, 4.0, 0.55, desc, 12, False, GRAY)


# ═══════════════════════════════════════════════════════
# Slide 8: 参考文献 (精简)
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
header(slide, "关键参考文献")

refs_cols = [
    ("CXL 技术基础", [
        '[1] Sun et al., "Demystifying CXL," MICRO\'23',
        '[2] Li et al., "Pond: CXL Memory Pooling," ASPLOS\'23',
        '[3] Maruf et al., "TPP," ASPLOS\'23',
        '[4] Gouk et al., "DirectCXL," ATC\'22',
    ], 0.6, BLUE),
    ("LLM 推理系统", [
        '[5] Kwon et al., "vLLM/PagedAttention," SOSP\'23',
        '[6] Zhong et al., "DistServe," OSDI\'24',
        '[7] Patel et al., "Splitwise," ISCA\'24',
        '[8] Qin et al., "Mooncake (月之暗面)," arXiv\'24',
        '[9] NVIDIA, "Dynamo," 2025',
    ], 4.5, GREEN),
    ("CXL + LLM (最新)", [
        '[10] Yoon et al., "TraCT," arXiv 2025.12',
        '[11] Yang et al., "Beluga," SIGMOD\'26',
        '[12] Gao et al., "CachedAttention," ASPLOS\'25',
        '[13] Sheng et al., "FlexGen," ICML\'23',
    ], 8.5, ACCENT),
]
for (col_title, refs, xp, clr) in refs_cols:
    pill(slide, xp, 1.1, 3.5, 0.35, clr, col_title, 12)
    for i, r in enumerate(refs):
        tb(slide, xp + 0.1, 1.6 + i * 0.42, 3.8, 0.4, r, 11, False, DARK)

# 扩展参考
rect(slide, 0.6, 3.6, 12.2, 2.8, RGBColor(0xF8, 0xF8, 0xF8), RGBColor(0xE0, 0xE0, 0xE0))
tb(slide, 0.8, 3.65, 11, 0.35, "扩展参考 (30+ 篇)", 13, True, GRAY)
ext = [
    "SGLang (arXiv'23) · FlashAttention-2 (ICLR'24) · Infinite-LLM (SOSP'24) · LoongServe (SOSP'24)",
    "HeteGen (MLSys'24) · vAttention (ASPLOS'25) · CacheBlend (EuroSys'25) · RAGCache (arXiv'24)",
    "MemServe (arXiv'24) · CacheGen (SIGCOMM'24) · InfiniGen (OSDI'24) · Orca (OSDI'22)",
    "MEMTIS (SOSP'23) · CXL-SHM (HPCA'24) · AWQ (MLSys'24) · GPTQ (ICLR'23) · DeepSpeed (SC'22)",
    "Sangam (arXiv'25) · TRACE (arXiv'25) · PIM Is All You Need (ASPLOS'25) · CXLRAMSim (arXiv'26)",
]
for i, line in enumerate(ext):
    tb(slide, 0.8, 4.05 + i * 0.45, 11.5, 0.4, line, 11, False, GRAY)


# ═══════════════════════════════════════════════════════
# Slide 9: 谢谢
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# 顶部色带
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

# 左侧装饰
deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.08), Inches(0.12), Inches(7.42))
deco.fill.solid(); deco.fill.fore_color.rgb = BLUE; deco.line.fill.background()
deco2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.12), Inches(0.08), Inches(0.05), Inches(7.42))
deco2.fill.solid(); deco2.fill.fore_color.rgb = ACCENT; deco2.line.fill.background()

tb(slide, 1.5, 2.5, 10, 1.0, "谢谢！", 52, True, BLUE, PP_ALIGN.CENTER)
tb(slide, 1.5, 3.8, 10, 0.6, "敬请各位老师批评指正", 22, False, GRAY, PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.8), Inches(3.5), Pt(2))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

tb(slide, 1.5, 5.2, 10, 0.5, "何宸禹    浙江大学计算机学院", 16, False, GRAY, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# 写入备注（发言稿）
# ═══════════════════════════════════════════════════════
speaker_notes = [
    # Slide 0: 封面
    "各位老师好，我是何宸禹，我们团队的课题是《面向大模型推理加速的CXL内存扩展与异构计算架构研究》。",

    # Slide 1: 大模型推理流程
    ("在进入正题之前，我先简单介绍一下大模型推理的工作流程。\n\n"
     "当用户输入一段文字后，推理分为两个阶段。第一个是Prefill阶段，GPU并行处理所有输入token，"
     "这是计算密集型的，同时会产生KV Cache——也就是每一层Transformer的Key和Value矩阵的缓存。\n\n"
     "第二个是Decode阶段，模型逐个生成输出token，每生成一个token都需要读取之前所有token的KV Cache，"
     "所以这个阶段是内存密集型的，而且KV Cache会随序列长度线性增长。\n\n"
     "右下方的条形图展示了70B模型在不同上下文长度下的KV Cache大小——"
     "128K上下文就需要近20GB，这还只是单个请求。多个并发请求下，GPU的80GB HBM很快就会被占满。"
     "这就是我们课题要解决的核心问题。"),

    # Slide 2: 内存墙
    ("再看更宏观的数据。左边这张图展示了GPU算力与内存增长的'剪刀差'："
     "从V100到B200，算力增长了14倍，但HBM容量仅增长6倍，差距持续扩大。\n\n"
     "右边是具体的内存压力数据。以70B模型加128K上下文为例，"
     "KV Cache就需要约160GB，是模型权重的2倍。而月之暗面的Kimi，50M token的KV Cache需要20TB。"
     "GPU的80GB HBM远远不够。"),

    # Slide 3: P/D解耦与CXL
    ("业界的主流趋势是将推理拆分为Prefill和Decode两个阶段独立部署，"
     "DistServe等工作已经证明这能提升7倍以上的吞吐。但解耦后带来新问题："
     "Prefill生成的KV Cache需要传输给Decode实例，目前依赖RDMA，延迟高、编程复杂。\n\n"
     "下方展示了内存层次，我们可以看到，在DRAM和SSD之间存在一个巨大的性能鸿沟。"
     "CXL技术恰好填补了这个空白——它提供TB级容量、约200ns的低延迟、load/store语义访问，而且支持多主机内存池化。"),

    # Slide 4: TraCT & Beluga
    ("这是我们要重点介绍的——2025年的两篇最新论文已经实证验证了CXL用于LLM推理的可行性。\n\n"
     "左边是Virginia Tech和SK Hynix合作的TraCT，它用CXL共享内存替代RDMA做P/D间KV传输，"
     "基于NVIDIA Dynamo框架，GPU通过DMA直接读写CXL内存，完全消除了NIC中转。"
     "结果非常显著：TTFT降低9.8倍，P99延迟降低6.2倍，吞吐提升1.6倍。\n\n"
     "右边是阿里巴巴的Beluga，已被SIGMOD'26录用。它用CXL 2.0 Switch构建8TB共享内存池来管理KV Cache，"
     "集成到vLLM后，TTFT降低89.6%，吞吐提升7.35倍。这也是首个让GPU直接访问CXL内存池的系统。\n\n"
     "这两篇论文恰好验证了我们课题的两个研究方向。"),

    # Slide 5: 两大方向总览
    ("基于以上分析，我们提出两个互补的研究方向。\n\n"
     "方向一：Serverless场景下P/D间KV Cache通讯。核心是用CXL共享内存池替代RDMA，"
     "实现P/D间零拷贝KV传递。KV Cache独立于实例生命周期，存活在CXL池中，支持弹性伸缩。\n\n"
     "方向二：KV Cache Offloading到CXL内存。构建HBM→DRAM→CXL→SSD四层架构，"
     "用CXL填补DRAM与SSD之间的性能鸿沟。热KV留在HBM，冷KV下沉到CXL。\n\n"
     "底部这句话概括了两者关系——方向一解决P/D间传输效率，方向二解决单实例容量瓶颈。"),

    # Slide 6: 方向一详细
    ("这页是方向一的核心对比，左右两个方案一目了然。\n\n"
     "左边是当前RDMA方案：数据从Prefill GPU到RDMA网卡再到对端，6次内存拷贝，"
     "75%的延迟来自同步开销，而且RDMA的QP管理编程非常复杂。\n\n"
     "右边是我们的CXL方案：Prefill GPU通过DMA直接写入CXL共享内存池，Decode GPU直接读取，"
     "采用load/store语义，零拷贝。中间蓝色条形图显示，CXL方案的延迟仅为RDMA的十分之一。\n\n"
     "下方三个研究子课题分别是：KV共享协议设计、Serverless弹性调度、以及性能优化与对比评测。"),

    # Slide 7: 方向二详细
    ("这页展示方向二的分层存储架构。左侧是四级存储层次：从最快的GPU HBM到CPU DRAM"
     "再到CXL内存最后到SSD，容量逐级增大、速度逐级降低。CXL层用蓝色高亮，是我们新增的关键一层。\n\n"
     "右上方红色卡片指出现有瓶颈：FlexGen的GPU与CPU间仅12GB/s，SSD延迟是DRAM的100倍，"
     "中间缺少过渡层。右下方绿色卡片展示CXL的优势——38GB/s带宽，是PCIe的3倍、SSD的5倍，完美填补空白。\n\n"
     "研究内容包括：CXL-aware热度感知分层策略、vLLM/SGLang的CXL后端集成、"
     "以及多场景benchmark评测，目标是吞吐提升30%以上。"),

    # Slide 8: 时间规划
    ("时间规划分四个阶段：前3个月文献调研和环境搭建，中间两个阶段分别推进两个方向的研究，"
     "最后2-3个月整合结题和论文撰写。\n\n"
     "预期贡献四个方面：P/D间CXL KV通讯方案、CXL KV Offloading系统、"
     "对接蚂蚁智算平台的实际系统验证与落地、以及开源系统原型。"),

    # Slide 9: 参考文献
    "这是我们的参考文献，包括CXL基础、LLM推理系统、以及最新的CXL+LLM交叉论文，共30余篇。",

    # Slide 10: 谢谢
    "以上是我们的开题报告，请各位老师批评指正，谢谢！",
]

for i, note_text in enumerate(speaker_notes):
    if i < len(prs.slides):
        notes_slide = prs.slides[i].notes_slide
        notes_slide.notes_text_frame.text = note_text

# ═══════════════════════════════════════════════════════
# 保存
# ═══════════════════════════════════════════════════════
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "开题报告_何宸禹_v5.pptx")
prs.save(out_path)
print(f"PPT已生成: {out_path}")
print(f"共 {len(prs.slides)} 页")
