#!/usr/bin/env python3
"""
generate_ppt.py — marginlens PPT export worker.

Reads pending jobs from Supabase `ppt_jobs`, generates a .pptx file that
visually matches the dark-themed interactive web presentation, uploads the
file to Supabase Storage (`ppt-exports` bucket), and updates job status.

Required env vars:
    SUPABASE_URL          e.g. https://xxxx.supabase.co
    SUPABASE_SERVICE_KEY  service_role key (bypasses RLS)

Dependencies:
    pip install python-pptx requests matplotlib

Run:
    python generate_ppt.py
"""

import io, json, os, re, sys
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette matching the web theme ────────────────────────────────────
BG          = RGBColor(0x08, 0x08, 0x0f)
FG          = RGBColor(0xf0, 0xf0, 0xf8)
FG_DIM      = RGBColor(0x70, 0x70, 0xa0)
FG_MID      = RGBColor(0xc0, 0xc0, 0xe0)
SURFACE     = RGBColor(0x0f, 0x0f, 0x22)
SURFACE2    = RGBColor(0x14, 0x14, 0x30)
BORDER      = RGBColor(0x25, 0x25, 0x55)

ACCENT: dict[str, RGBColor] = {
    'blue':   RGBColor(0x5b, 0x9c, 0xf8),
    'purple': RGBColor(0xb0, 0x6a, 0xf4),
    'green':  RGBColor(0x4a, 0xde, 0x80),
    'amber':  RGBColor(0xfb, 0xbf, 0x24),
    'rose':   RGBColor(0xfb, 0x71, 0x85),
    'gray':   RGBColor(0x94, 0xa3, 0xb8),
}
CALLOUT_COLORS: dict[str, tuple[RGBColor, str]] = {
    'note':     (ACCENT['blue'],   '📝'),
    'tip':      (ACCENT['green'],  '💡'),
    'warning':  (ACCENT['amber'],  '⚠'),
    'question': (ACCENT['purple'], '?'),
    'insight':  (ACCENT['rose'],   '✦'),
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = 'Segoe UI'   # Embeds as font name; PowerPoint subs if unavailable


def get_accent(m: dict) -> RGBColor:
    return ACCENT.get(m.get('accent', ''), ACCENT['blue'])


# ── Markdown → plain text (strip markers, keep content) ─────────────────────
def strip_md(text: str) -> str:
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    text = re.sub(r'`(.*?)`', r'\1', text)
    text = re.sub(r'#{1,6}\s+', '', text)
    text = re.sub(r'!\[.*?\]\(.*?\)', '', text)
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    return text.strip()


# ── Low-level drawing helpers ────────────────────────────────────────────────
def dark_bg(slide) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG


def accent_bar(slide, color: RGBColor, height_in: float = 0.055) -> None:
    """Thin color bar at top of slide."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(height_in))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def add_rect(slide, l, t, w, h,
             fill: RGBColor = SURFACE,
             border: RGBColor | None = BORDER,
             border_pt: float = 1.0,
             rounded: bool = True) -> object:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    return s


def add_text(slide, l, t, w, h, text: str,
             size: float = 16,
             bold: bool = False,
             italic: bool = False,
             color: RGBColor = FG,
             align: PP_ALIGN = PP_ALIGN.LEFT) -> object:
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    return tf


def add_para(tf, text: str,
             size: float = 14,
             bold: bool = False,
             color: RGBColor = FG,
             align: PP_ALIGN = PP_ALIGN.LEFT,
             space_before: float = 0) -> object:
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    return p


def slide_title(slide, title: str, color: RGBColor = FG) -> None:
    add_text(slide, 1.2, 0.18, 11.0, 0.55, title, size=22, bold=True, color=color)


# ── Module renderers ─────────────────────────────────────────────────────────
def render_hero(prs: Presentation, m: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    color = get_accent(m)
    accent_bar(slide, color, 0.08)

    y = 1.4
    chips = m.get('chips') or []
    if chips:
        x = 1.6
        for chip in chips[:6]:
            chip_w = max(len(chip) * 0.115 + 0.45, 0.9)
            add_rect(slide, x, y, chip_w, 0.33, RGBColor(0x10, 0x14, 0x2c), color)
            add_text(slide, x + 0.06, y + 0.03, chip_w - 0.1, 0.27,
                     chip, 11, color=color, align=PP_ALIGN.CENTER)
            x += chip_w + 0.18
        y = 2.05

    add_text(slide, 1.6, y, 10.5, 1.6, m.get('title', ''), 44, bold=True,
             color=FG, align=PP_ALIGN.CENTER)
    y += 1.65

    if m.get('subtitle'):
        add_text(slide, 1.6, y, 10.5, 0.65, m['subtitle'], 22,
                 color=FG_DIM, align=PP_ALIGN.CENTER)
        y += 0.75

    if m.get('summary'):
        add_text(slide, 2.2, y, 9.5, 1.4, strip_md(m['summary']), 16,
                 color=RGBColor(0x90, 0x90, 0xb8), align=PP_ALIGN.CENTER)


def render_section(prs: Presentation, m: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    color = get_accent(m)
    accent_bar(slide, color)

    # Left accent stripe
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(0.9), Inches(0.055), Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, 1.4, 0.9, 11.0, 0.72, m.get('title', ''), 28, bold=True, color=FG)

    content = strip_md(m.get('content', ''))
    add_text(slide, 1.2, 1.9, 11.4, 5.0, content, 17, color=FG_MID)


def render_keypoints(prs: Presentation, m: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    color = get_accent(m)
    accent_bar(slide, color)

    slide_title(slide, m.get('title') or '关键要点', color)

    items = m.get('items', [])
    for i, item in enumerate(items[:8]):
        y = 0.95 + i * 0.73
        # Colored bullet dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(1.2), Inches(y + 0.14),
                                     Inches(0.14), Inches(0.14))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_text(slide, 1.48, y, 11.5, 0.65, strip_md(item), 17, color=FG)


def render_definition(prs: Presentation, m: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    color = get_accent(m)
    accent_bar(slide, color)

    add_text(slide, 1.2, 0.3, 11.0, 0.65, m.get('term', ''), 32, bold=True, color=color)

    # Divider
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.1), Inches(11.3), Pt(0.8))
    div.fill.solid()
    div.fill.fore_color.rgb = BORDER
    div.line.fill.background()

    defn = strip_md(m.get('definition', ''))
    add_text(slide, 1.2, 1.25, 11.3, 3.8, defn, 18, color=FG_MID)

    if m.get('example'):
        add_rect(slide, 1.2, 5.1, 11.3, 1.7, SURFACE2, color)
        add_text(slide, 1.45, 5.18, 1.2, 0.38, '示例', 13, bold=True, color=color)
        add_text(slide, 1.45, 5.58, 11.0, 0.95, strip_md(m['example']), 16, color=FG)


def render_formula(prs: Presentation, m: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    color = get_accent(m)
    accent_bar(slide, color)

    if m.get('caption'):
        add_text(slide, 1.5, 0.3, 11.0, 0.6, m['caption'], 22, bold=True, color=FG)

    latex = m.get('latex', '')
    img_inserted = False

    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt

        fig = plt.figure(figsize=(9, 2.5), facecolor='#0d0d1a')
        ax = fig.add_axes([0, 0, 1, 1])
        ax.set_facecolor('#0d0d1a')
        ax.axis('off')
        ax.text(0.5, 0.5, f'${latex}$', transform=ax.transAxes,
                fontsize=30, ha='center', va='center', color='#f0f0f8')
        buf = io.BytesIO()
        plt.savefig(buf, format='png', bbox_inches='tight',
                    facecolor='#0d0d1a', edgecolor='none', dpi=150)
        plt.close(fig)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(1.5), Inches(1.1), Inches(10.5), Inches(2.8))
        img_inserted = True
    except Exception as exc:
        print(f'  latex render failed: {exc}', file=sys.stderr)

    if not img_inserted:
        add_rect(slide, 1.5, 1.1, 10.5, 2.2, SURFACE2, color)
        add_text(slide, 1.7, 1.3, 10.1, 1.8, latex, 20,
                 color=color, align=PP_ALIGN.CENTER)

    if m.get('explanation'):
        add_text(slide, 1.5, 4.2, 10.5, 2.4,
                 strip_md(m['explanation']), 16, color=FG_DIM)


def render_callout(prs: Presentation, m: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    variant = m.get('variant', 'note')
    color, icon = CALLOUT_COLORS.get(variant, CALLOUT_COLORS['note'])
    accent_bar(slide, color)

    add_rect(slide, 1.15, 0.8, 11.3, 5.8, SURFACE2, color)

    header_text = f'{icon}  {m.get("title") or variant.title()}'
    add_text(slide, 1.45, 0.95, 10.8, 0.55, header_text, 20, bold=True, color=color)

    body = strip_md(m.get('body', ''))
    add_text(slide, 1.45, 1.7, 10.8, 4.5, body, 18, color=FG)


def render_qa(prs: Presentation, m: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    color = get_accent(m)
    accent_bar(slide, color)

    # Q block
    add_text(slide, 1.2, 0.25, 0.55, 0.65, 'Q', 34, bold=True, color=color)
    add_text(slide, 1.85, 0.3, 11.0, 0.88, m.get('question', ''), 22, bold=True, color=FG)

    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.35), Inches(11.3), Pt(0.8))
    div.fill.solid()
    div.fill.fore_color.rgb = BORDER
    div.line.fill.background()

    # A block
    add_text(slide, 1.2, 1.55, 0.55, 0.65, 'A', 34, bold=True,
             color=ACCENT['green'])
    answer = strip_md(m.get('answer', ''))
    add_text(slide, 1.85, 1.6, 11.0, 5.0, answer, 17,
             color=RGBColor(0xd0, 0xee, 0xd0))


def render_quiz(prs: Presentation, m: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    color = get_accent(m)
    accent_bar(slide, color)

    add_text(slide, 1.2, 0.25, 11.5, 0.95, m.get('question', ''), 22, bold=True, color=FG)

    correct = m.get('correctIndex', 0)
    labels = ['A', 'B', 'C', 'D', 'E']
    options = m.get('options', [])

    for i, option in enumerate(options[:5]):
        y = 1.38 + i * 0.88
        is_ok = i == correct
        bg = RGBColor(0x0c, 0x1e, 0x0c) if is_ok else SURFACE
        bd = ACCENT['green'] if is_ok else BORDER
        add_rect(slide, 1.2, y, 11.5, 0.72, bg, bd)
        lbl_color = ACCENT['green'] if is_ok else color
        add_text(slide, 1.38, y + 0.13, 0.48, 0.46,
                 labels[i], 17, bold=True, color=lbl_color)
        add_text(slide, 1.95, y + 0.13, 10.6, 0.46,
                 option, 16, color=RGBColor(0xc0, 0xee, 0xc0) if is_ok else FG)

    if m.get('explanation'):
        expl = strip_md(m['explanation'])
        add_text(slide, 1.2, 6.55, 11.5, 0.78,
                 f'💡 {expl}', 13, color=ACCENT['green'])


def render_summary(prs: Presentation, m: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    color = get_accent(m)
    accent_bar(slide, color)

    slide_title(slide, m.get('title') or '总结', color)

    points = m.get('points', [])
    for i, pt in enumerate(points[:7]):
        y = 0.95 + i * 0.75
        # Numbered badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(y),
            Inches(0.34), Inches(0.34))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BG
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
        add_text(slide, 1.68, y + 0.02, 11.5, 0.6, strip_md(pt), 17, color=FG)


MODULE_RENDERERS = {
    'hero':       render_hero,
    'section':    render_section,
    'keypoints':  render_keypoints,
    'definition': render_definition,
    'formula':    render_formula,
    'callout':    render_callout,
    'qa':         render_qa,
    'quiz':       render_quiz,
    'summary':    render_summary,
}


# ── Core generation ──────────────────────────────────────────────────────────
def generate_pptx(site_data: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    modules = site_data.get('modules', [])
    if not modules:
        raise ValueError('site_data contains no modules')

    for m in modules:
        mtype = m.get('type', '')
        renderer = MODULE_RENDERERS.get(mtype)
        if renderer:
            try:
                renderer(prs, m)
            except Exception as exc:
                print(f'  [warn] module {m.get("id")} ({mtype}): {exc}', file=sys.stderr)
        else:
            print(f'  [skip] unknown module type: {mtype}', file=sys.stderr)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ── Supabase I/O ─────────────────────────────────────────────────────────────
def supabase_headers(key: str) -> dict:
    return {
        'apikey': key,
        'Authorization': f'Bearer {key}',
        'Content-Type': 'application/json',
    }


def fetch_pending_jobs(url: str, key: str) -> list[dict]:
    r = requests.get(
        f'{url}/rest/v1/ppt_jobs',
        params={'status': 'eq.pending', 'limit': '10', 'order': 'created_at.asc'},
        headers=supabase_headers(key),
        timeout=15,
    )
    r.raise_for_status()
    return r.json()


def set_job_status(url: str, key: str, job_id: str, **fields) -> None:
    r = requests.patch(
        f'{url}/rest/v1/ppt_jobs',
        params={'id': f'eq.{job_id}'},
        headers=supabase_headers(key),
        json=fields,
        timeout=15,
    )
    r.raise_for_status()


def upload_pptx(url: str, key: str, filename: str, data: bytes) -> str:
    upload_headers = {
        'apikey': key,
        'Authorization': f'Bearer {key}',
        'Content-Type': (
            'application/vnd.openxmlformats-officedocument'
            '.presentationml.presentation'
        ),
    }
    r = requests.post(
        f'{url}/storage/v1/object/ppt-exports/{filename}',
        headers=upload_headers,
        data=data,
        timeout=60,
    )
    if r.status_code not in (200, 201):
        raise RuntimeError(f'Storage upload failed [{r.status_code}]: {r.text}')
    return f'{url}/storage/v1/object/public/ppt-exports/{filename}'


# ── Entry point ──────────────────────────────────────────────────────────────
def main() -> None:
    sb_url = os.environ['SUPABASE_URL'].rstrip('/')
    sb_key = os.environ['SUPABASE_SERVICE_KEY']

    jobs = fetch_pending_jobs(sb_url, sb_key)
    if not jobs:
        print('No pending PPT jobs.')
        return

    print(f'Processing {len(jobs)} job(s)…')
    for job in jobs:
        job_id = job['id']
        print(f'  [{job_id}] starting…')

        # Mark as processing so concurrent runners skip it
        set_job_status(sb_url, sb_key, job_id, status='processing')

        try:
            pptx_bytes = generate_pptx(job['site_data'])
            filename = f'{job_id}.pptx'
            pptx_url = upload_pptx(sb_url, sb_key, filename, pptx_bytes)
            set_job_status(sb_url, sb_key, job_id, status='done', pptx_url=pptx_url)
            print(f'  [{job_id}] done → {pptx_url}')
        except Exception as exc:
            msg = str(exc)
            set_job_status(sb_url, sb_key, job_id, status='error', error_msg=msg)
            print(f'  [{job_id}] error: {msg}', file=sys.stderr)


if __name__ == '__main__':
    main()
